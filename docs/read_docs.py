import os
import glob
try:
    from pypdf import PdfReader
except ImportError:
    print("Could not import pypdf")
try:
    from pptx import Presentation
except ImportError:
    print("Could not import pptx")

DOCS_DIR = "/media/marco/Dados/projeto_pdi/docs"
OUTPUT_FILE = "/media/marco/Dados/projeto_pdi/docs/dump.txt"

def extract_pdf(file_path):
    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            t = page.extract_text()
            if t: text += t + "\n"
    except Exception as e:
        text = f"Error reading PDF: {e}"
    return text

def extract_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Error reading PPTX: {e}"
    return text

with open(OUTPUT_FILE, "w", encoding="utf-8") as f_out:
    for file in sorted(glob.glob(os.path.join(DOCS_DIR, "*.pdf"))):
        f_out.write(f"\n{'='*50}\nFILE: {os.path.basename(file)}\n{'='*50}\n")
        f_out.write(extract_pdf(file)[:2000]) # first 2000 chars for preview
        f_out.write("\n\n")

    for file in sorted(glob.glob(os.path.join(DOCS_DIR, "*.pptx"))):
        f_out.write(f"\n{'='*50}\nFILE: {os.path.basename(file)}\n{'='*50}\n")
        f_out.write(extract_pptx(file)[:2000])
        f_out.write("\n\n")

print("Dump complete!")
